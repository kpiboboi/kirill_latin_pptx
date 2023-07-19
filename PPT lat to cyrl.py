from pptx import Presentation
from transliterate import translit

def reverse_transliterate_presentation(presentation_file):
    prs = Presentation(presentation_file)
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        latin_text = run.text
                        cyrillic_text = translit(latin_text, 'ru')
                        run.text = cyrillic_text
    
    new_presentation_file = r'C:\Users\s.ibodov\Downloads\tst ppt\Bosh ofis Tashkiliy tuzilma.pptx'
    prs.save(new_presentation_file)

# Пример использования:
presentation_file = r'C:\Users\s.ibodov\Downloads\tst ppt\Bosh ofis Tashkiliy tuzilma.pptx'
reverse_transliterate_presentation(presentation_file)
print("Done")
