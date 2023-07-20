from pptx import Presentation
from transliterate import translit
import re
import os

def transliterate_word(word):
    translit_rules = {
        "ю" : "yu", "Ю" : "Yu",
        "ў" : "o'", "Ў" : "O'",
        "ё" : "yo", "Ё" : "Yo",
        "ғ" : "g'", "Ғ" : "G'",
        "қ" : "q", "Қ" : "Q",
        "ҳ" : "h", "Ҳ" : "H",
        "х" : "x", "Х" : "X",
        "ж" : "j", "Ж" : "J",
        "й" : "y", "Й" : "Y",
        "я" : "ya", "Я" : "Ya",
        "ц" : "s", "Ц" : "S",
        "ы" : "i", 
        # Agar kerak bo'lsa, boshqa almashtirishlarni qo'shishingiz mumkin
    }

    # So'zdagi belgilarni qoidalarga muvofiq almashtiramiz
    for cyrillic_char, latin_char in translit_rules.items():
        word = word.replace(cyrillic_char, latin_char)

    # So'z boshida "e" ni va "ye" bilan almashtiramiz
    if word.startswith("е"):
        if len(word) > 1:
            word = "ye" + word[1:]
        else:
            word = "ye"
    elif word.startswith("Е"):
        if len(word) > 1:
            word = "Ye" + word[1:]
        else:
            word = "Ye"
    return word
pass

def transliterate_presentation(presentation_file):
    prs = Presentation(presentation_file)

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        cyrillic_text = run.text
                        words = re.split(r"(\s+)", cyrillic_text)
                        latin_words = [transliterate_word(word) for word in words]
                        latin_text = "".join(latin_words)
                        run.text = translit(latin_text, "ru", reversed=True)

    prs.save(presentation_file)

def process_presentation_file(file_path):
    try:
        transliterate_presentation(file_path)
        print(f"✅ Fayl muvafaqqiyatli o'zgartirildi: {file_path}")
        return True  # Возвращаем True в случае успешного выполнения
    except Exception as e:
        print(f"❌ Fayl o'zgartirishida xatolik: {file_path}. Sabab: {str(e)}")
        return False  # Возвращаем False в случае ошибки

def process_folder(folder_path):
    for root, dirs, files in os.walk(folder_path):
        for file in files:
            if file.endswith('.pptx'):
                file_path = os.path.join(root, file)
                success = process_presentation_file(file_path)
                if success:
                    print(f"✅ Fayl muvafaqqiyatli o'zgartirildi: {file_path}")
                else:
                    print(f"❌ Fayl o'zgartirishida xatolik: {file_path}")

# Natija:
presentation_folder = r"C:\Users\s.ibodov\Downloads\tst ppt" #Papka manzilini kiriting
process_folder(presentation_folder)
print("✅✅✅ PPTX fayl kirilchadan lotinchaga muvafaqqiyatli o'girildi ✅✅✅")